from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
MANIFEST_PATH = ROOT / "control/project_domain/resources/pptx_jobs/02_1/manifest.json"
BUNDLE_PATH = ROOT / "control/project_domain/resources/manifests/ppt_regeneration_handoff_bundle_v0_1_at2026_04_13.json"
OUTPUT_DIR = ROOT / "control/project_domain/resources/assets/portfolio_drafts/lean_02_1_system_first_v2"
RENDER_SOURCE_DIR = OUTPUT_DIR / "render_sources"
ROLE_MATRIX_PATH = ROOT / "control/project_domain/resources/manifests/lean_02_1_system_first_v2_image_role_matrix_at2026_04_13.json"
DECK_PATH = OUTPUT_DIR / "lean_02_1_system_first_v2.pptx"
AUTHORING_SKILL = "<CODEX_HOME>/skills/pptx/SKILL.md"
COPY_CLARITY_SKILL = "<CLAUDE_SKILLS_ROOT>/semantic-clarity-enhanced/SKILL.md"

SLIDE_SIZE = (13.333, 7.5)

DARK_BG = RGBColor(0x18, 0x22, 0x2B)
LIGHT_BG = RGBColor(0xF7, 0xF2, 0xEA)
PANEL_BG = RGBColor(0xEA, 0xE2, 0xD5)
DEEP_TEAL = RGBColor(0x2D, 0x8C, 0x88)
WARM_SAND = RGBColor(0xD7, 0xBF, 0x9D)
DARK_TEXT = RGBColor(0x1E, 0x24, 0x2C)
LIGHT_TEXT = RGBColor(0xFA, 0xF7, 0xF2)
MUTED_TEXT = RGBColor(0x64, 0x69, 0x6F)
FONT_NAME = "Apple SD Gothic Neo"


@dataclass(frozen=True)
class SlideSpec:
    slide_no: int
    title: str
    subtitle: str
    primary_image: str
    support_images: tuple[str, ...]
    presentation_role: str
    visual_type: str
    supporting_text_goal: str
    layout_role: str
    footer: str
    top_bottom: bool
    notes: tuple[str, ...]


SLIDE_SPECS = [
    SlideSpec(
        slide_no=1,
        title="시스템을 구조로 설계하고 증명합니다",
        subtitle="아키텍처 다이어그램을 크게 두고, 아래에서 역할만 짧게 해석합니다.",
        primary_image="image27.png",
        support_images=(),
        presentation_role="architecture_hero",
        visual_type="system_diagram",
        supporting_text_goal="시스템 사고가 첫 인상으로 보이도록 만든다.",
        layout_role="top_header_large_diagram_bottom_notes",
        footer="source: 02_1/image27.png • architecture hero • regenerated v2",
        top_bottom=True,
        notes=(
            "이미지 자체가 시스템 설명의 본문이다.",
            "검색·생성·검증 흐름을 한 화면에서 읽게 한다.",
        ),
    ),
    SlideSpec(
        slide_no=2,
        title="입력과 결과가 한 흐름으로 이어집니다",
        subtitle="설정 UI와 생성 결과 UI를 한 장의 경험으로 연결합니다.",
        primary_image="image4.png",
        support_images=("image37.png",),
        presentation_role="product_configuration_ui",
        visual_type="ui_screenshot_pair",
        supporting_text_goal="설정 화면과 결과 화면이 하나의 경험으로 이어진다는 점을 보여준다.",
        layout_role="top_header_pair_ui_bottom_notes",
        footer="source: 02_1/image4.png + image37.png • paired product flow • regenerated v2",
        top_bottom=True,
        notes=(
            "왼쪽은 조건을 입력하는 화면, 오른쪽은 결과를 확인하는 화면이다.",
            "텍스트보다 화면 관계가 먼저 읽히도록 배치한다.",
        ),
    ),
    SlideSpec(
        slide_no=3,
        title="개인 맥락이 신뢰를 만듭니다",
        subtitle="프로필 요약과 인물 이미지를 분리하지 않고 한 장에 묶습니다.",
        primary_image="image5.png",
        support_images=("image6.jpeg",),
        presentation_role="profile_summary",
        visual_type="profile_summary_plus_portrait",
        supporting_text_goal="정리된 이력과 실제 인물 이미지가 함께 신뢰도를 만든다는 점을 강조한다.",
        layout_role="top_header_profile_band_bottom_notes",
        footer="source: 02_1/image5.png + image6.jpeg • profile credibility pair • regenerated v2",
        top_bottom=True,
        notes=(
            "이력 요약이 주인공이고, portrait는 신뢰를 보강한다.",
            "프로필을 글이 아니라 실제 자료로 보여준다.",
        ),
    ),
    SlideSpec(
        slide_no=4,
        title="문제와 구현을 한 화면에 둡니다",
        subtitle="문제 설명과 Python 구현이 동시에 보여야 해결 증거가 됩니다.",
        primary_image="image22.png",
        support_images=(),
        presentation_role="problem_and_code_proof",
        visual_type="code_and_problem_screen",
        supporting_text_goal="문제 해결형 개발자라는 점을 시각 증거로 보여준다.",
        layout_role="top_header_full_problem_code_bottom_notes",
        footer="source: 02_1/image22.png • problem and code proof • regenerated v2",
        top_bottom=True,
        notes=(
            "문제 정의와 코드가 함께 남아야 구현 깊이가 보인다.",
        ),
    ),
    SlideSpec(
        slide_no=5,
        title="프로젝트 표가 실행 이력을 증명합니다",
        subtitle="표는 장식이 아니라 evidence이므로 값 읽기 가능성을 최대한 보존합니다.",
        primary_image="image23.png",
        support_images=(),
        presentation_role="portfolio_evidence_table",
        visual_type="table_evidence",
        supporting_text_goal="실행 범위와 지속성을 포트폴리오 증거 표로 보여준다.",
        layout_role="top_header_full_table_bottom_notes",
        footer="source: 02_1/image23.png • portfolio evidence table • regenerated v2",
        top_bottom=True,
        notes=(
            "표 내부 값이 읽혀야 이 slide가 의미를 가진다.",
        ),
    ),
    SlideSpec(
        slide_no=6,
        title="적용형 AI 흐름을 끝까지 닫습니다",
        subtitle="업로드, 조건 선택, 결과 화면을 하나의 applied workflow로 읽히게 합니다.",
        primary_image="image29.png",
        support_images=(),
        presentation_role="applied_ai_workflow_ui",
        visual_type="workflow_ui",
        supporting_text_goal="실험 결과를 실제 사용 흐름으로 이어 붙일 수 있다는 점을 강조한다.",
        layout_role="top_header_tall_ui_bottom_notes",
        footer="source: 02_1/image29.png • applied AI workflow UI • regenerated v2",
        top_bottom=True,
        notes=(
            "멀티모달 입력은 noise가 아니라 workflow context다.",
            "실험 결과를 실제 사용 흐름으로 연결하는 장면이다.",
        ),
    ),
]


def load_manifest_index() -> dict[str, dict]:
    data = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    return {item["file"]: item for item in data["exported_images"]}


def load_bundle() -> dict:
    return json.loads(BUNDLE_PATH.read_text(encoding="utf-8"))


def rgb_to_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def repo_relative(path: Path) -> str:
    return str(path.relative_to(ROOT))


def add_rect(slide, x, y, w, h, fill, line=None, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_tag(slide, x, y, text, *, bg, fg):
    width = max(1.2, 0.13 * len(text) + 0.35)
    add_rect(slide, x, y, width, 0.42, bg)
    add_textbox(slide, x + 0.08, y + 0.02, width - 0.16, 0.32, text, size=10.5, color=fg, bold=True, valign=MSO_ANCHOR.MIDDLE)
    return width


def add_picture_contained(slide, image_path: Path, x, y, w, h, *, frame_fill=None, frame_line=None):
    if frame_fill is not None:
        add_rect(slide, x, y, w, h, frame_fill, line=frame_line)
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    image_ratio = img_w / img_h
    box_ratio = w / h
    if image_ratio > box_ratio:
        target_w = w
        target_h = w / image_ratio
    else:
        target_h = h
        target_w = h * image_ratio
    target_x = x + (w - target_w) / 2
    target_y = y + (h - target_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(target_x), Inches(target_y), width=Inches(target_w), height=Inches(target_h))


def base_slide(prs: Presentation, *, dark: bool):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG if dark else LIGHT_BG
    return slide


def add_footer(slide, text, *, dark):
    add_textbox(
        slide,
        0.72,
        7.03,
        11.9,
        0.2,
        text,
        size=9.5,
        color=RGBColor(0xA7, 0xB0, 0xB6) if dark else RGBColor(0x7A, 0x7E, 0x84),
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_header(slide, spec: SlideSpec, *, dark: bool, accent_text: str):
    title_color = LIGHT_TEXT if dark else DARK_TEXT
    subtitle_color = RGBColor(0xD8, 0xE1, 0xE8) if dark else MUTED_TEXT
    add_tag(slide, 0.78, 0.42, accent_text, bg=DEEP_TEAL, fg=DARK_BG)
    add_textbox(slide, 0.78, 0.92, 10.6, 0.56, spec.title, size=25.5 if spec.slide_no != 1 else 27, color=title_color, bold=True)
    add_textbox(slide, 0.8, 1.42, 11.5, 0.42, spec.subtitle, size=12.8, color=subtitle_color)


def add_note_panel(slide, x, y, w, h, notes, *, dark):
    bg = RGBColor(0x23, 0x30, 0x3B) if dark else PANEL_BG
    fg = LIGHT_TEXT if dark else DARK_TEXT
    add_rect(slide, x, y, w, h, bg, line=None)
    for idx, note in enumerate(notes):
        add_textbox(
            slide,
            x + 0.18,
            y + 0.14 + idx * 0.34,
            w - 0.36,
            0.26,
            f"• {note}",
            size=11.4,
            color=fg,
            valign=MSO_ANCHOR.MIDDLE,
        )


def render_slide(prs: Presentation, spec: SlideSpec, image_index: dict[str, dict]):
    dark = spec.slide_no in {1, 6}
    slide = base_slide(prs, dark=dark)
    primary_image_path = Path(image_index[spec.primary_image]["output_path"])
    support_paths = [Path(image_index[name]["output_path"]) for name in spec.support_images]
    add_header(slide, spec, dark=dark, accent_text=spec.presentation_role.replace("_", " ").upper())

    frame_fill = RGBColor(0x22, 0x2D, 0x36) if dark else RGBColor(0xEE, 0xE7, 0xDA)
    frame_line = RGBColor(0x3A, 0x48, 0x53) if dark else RGBColor(0xD7, 0xC8, 0xB2)

    if spec.slide_no == 1:
        add_picture_contained(slide, primary_image_path, 0.8, 1.95, 11.78, 4.5, frame_fill=frame_fill, frame_line=frame_line)
        add_note_panel(slide, 0.8, 6.05, 11.78, 0.74, spec.notes, dark=dark)
    elif spec.slide_no == 2:
        add_picture_contained(slide, primary_image_path, 0.82, 2.0, 5.65, 4.45, frame_fill=frame_fill, frame_line=frame_line)
        add_picture_contained(slide, support_paths[0], 6.72, 2.0, 5.88, 4.45, frame_fill=frame_fill, frame_line=frame_line)
        add_note_panel(slide, 0.82, 6.05, 11.78, 0.74, spec.notes, dark=dark)
    elif spec.slide_no == 3:
        add_picture_contained(slide, primary_image_path, 0.82, 1.95, 8.88, 4.55, frame_fill=frame_fill, frame_line=frame_line)
        add_picture_contained(slide, support_paths[0], 9.96, 2.18, 2.62, 3.9, frame_fill=frame_fill, frame_line=frame_line)
        add_note_panel(slide, 0.82, 6.08, 11.78, 0.68, spec.notes, dark=dark)
    elif spec.slide_no == 4:
        add_picture_contained(slide, primary_image_path, 0.82, 1.95, 11.78, 4.7, frame_fill=frame_fill, frame_line=frame_line)
        add_note_panel(slide, 0.82, 6.17, 11.78, 0.56, spec.notes, dark=dark)
    elif spec.slide_no == 5:
        add_picture_contained(slide, primary_image_path, 0.48, 1.82, 12.35, 5.12, frame_fill=frame_fill, frame_line=frame_line)
        add_note_panel(slide, 0.82, 6.36, 11.78, 0.34, spec.notes, dark=dark)
    else:
        add_picture_contained(slide, primary_image_path, 1.2, 1.86, 10.9, 5.12, frame_fill=frame_fill, frame_line=frame_line)
        add_note_panel(slide, 0.82, 6.26, 11.78, 0.46, spec.notes, dark=dark)

    add_footer(slide, spec.footer, dark=dark)


def build_role_matrix(image_index: dict[str, dict], bundle: dict) -> list[dict]:
    slide_map = {row["target_slide_no"]: row for row in bundle["slides"]}
    rows = []
    for spec in SLIDE_SPECS:
        slide_row = slide_map[spec.slide_no]
        rows.append(
            {
                "slide_no": spec.slide_no,
                "slide_title": spec.title,
                "visual_type": spec.visual_type,
                "image_filename": spec.primary_image,
                "support_images": list(spec.support_images),
                "portfolio_role": spec.presentation_role,
                "source_slide_numbers": image_index[spec.primary_image]["slide_usages"],
                "supporting_text_goal": spec.supporting_text_goal,
                "layout_role": spec.layout_role,
                "top_bottom_split": spec.top_bottom,
                "bundle_status": bundle["bundle_status"],
                "bundle_manual_remaining": slide_row["regeneration_handoff"]["manual_remaining"],
            }
        )
    return rows


def build_presentation(slides: list[SlideSpec], image_index: dict[str, dict], output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_SIZE[0])
    prs.slide_height = Inches(SLIDE_SIZE[1])
    for spec in slides:
        render_slide(prs, spec, image_index)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main():
    image_index = load_manifest_index()
    bundle = load_bundle()
    build_presentation(SLIDE_SPECS, image_index, DECK_PATH)
    RENDER_SOURCE_DIR.mkdir(parents=True, exist_ok=True)
    for spec in SLIDE_SPECS:
        build_presentation([spec], image_index, RENDER_SOURCE_DIR / f"slide-{spec.slide_no:02d}-source.pptx")
    ROLE_MATRIX_PATH.write_text(
        json.dumps(build_role_matrix(image_index, bundle), ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    print(
        json.dumps(
            {
                "deck": repo_relative(DECK_PATH),
                "render_sources": [repo_relative(RENDER_SOURCE_DIR / f"slide-{spec.slide_no:02d}-source.pptx") for spec in SLIDE_SPECS],
                "role_matrix": repo_relative(ROLE_MATRIX_PATH),
                "source_bundle": repo_relative(BUNDLE_PATH),
                "authoring_skill": AUTHORING_SKILL,
                "copy_clarity_skill": COPY_CLARITY_SKILL,
                "palette": {
                    "dark_bg": rgb_to_hex(DARK_BG),
                    "light_bg": rgb_to_hex(LIGHT_BG),
                    "accent_teal": rgb_to_hex(DEEP_TEAL),
                    "accent_sand": rgb_to_hex(WARM_SAND),
                },
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
