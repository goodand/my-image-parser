from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
MANIFEST_PATH = ROOT / "control/project_domain/resources/pptx_jobs/02_1/manifest.json"
OUTPUT_DIR = ROOT / "control/project_domain/resources/assets/portfolio_drafts/lean_02_1_system_first_v1"
RENDER_SOURCE_DIR = OUTPUT_DIR / "render_sources"
ROLE_MATRIX_PATH = ROOT / "control/project_domain/resources/manifests/lean_02_1_system_first_v1_image_role_matrix_at2026_04_11.json"
DECK_PATH = OUTPUT_DIR / "lean_02_1_system_first_v1.pptx"

SLIDE_SIZE = (13.333, 7.5)

DARK_BG = RGBColor(0x1C, 0x26, 0x31)
LIGHT_BG = RGBColor(0xF5, 0xF1, 0xEA)
MUTED_BG = RGBColor(0xE7, 0xDF, 0xD2)
TEAL = RGBColor(0x57, 0xB7, 0xA7)
SAND = RGBColor(0xD8, 0xC5, 0xAD)
DARK_TEXT = RGBColor(0x1F, 0x23, 0x29)
LIGHT_TEXT = RGBColor(0xFA, 0xF7, 0xF2)
SOFT_TEXT = RGBColor(0x8A, 0x95, 0x9E)
FONT_NAME = "Apple SD Gothic Neo"


@dataclass(frozen=True)
class SlideSpec:
    slide_no: int
    title: str
    subtitle: str
    primary_image: str
    support_images: tuple[str, ...]
    primary_role: str
    visual_type: str
    supporting_text_goal: str
    layout_role: str
    crop_required: bool
    bullets: tuple[str, ...]


SLIDE_SPECS = [
    SlideSpec(
        slide_no=1,
        title="시스템을 구조로 설명하는 개발자",
        subtitle="멀티모달 문서를 파싱하고 검색·생성 흐름을 설계한 결과를 다시 포트폴리오 자산으로 묶습니다.",
        primary_image="image27.png",
        support_images=(),
        primary_role="architecture_hero",
        visual_type="system_diagram",
        supporting_text_goal="시스템 사고가 첫 인상으로 보이도록 만든다.",
        layout_role="right_hero_diagram_with_left_title_stack",
        crop_required=False,
        bullets=(
            "이미지 입력을 단순 부속물이 아니라 시스템 단위로 다룹니다.",
            "UI, 검색, 생성, 결과 검증까지 한 화면 안에서 연결합니다.",
            "PPT 안의 이미지를 다시 증거 자료로 쓰는 흐름을 설계합니다.",
        ),
    ),
    SlideSpec(
        slide_no=2,
        title="입력과 결과로 증명하는 제품 경험",
        subtitle="입력 조건과 narrative 결과를 같은 사용자 흐름으로 읽히게 합니다.",
        primary_image="image4.png",
        support_images=("image37.png",),
        primary_role="product_configuration_ui",
        visual_type="ui_screenshot_pair",
        supporting_text_goal="설정 화면과 결과 화면이 하나의 경험으로 이어진다는 점을 보여준다.",
        layout_role="paired_ui_block_with_right_explanation",
        crop_required=False,
        bullets=(
            "입력 화면에서 플랫폼, 감정, 관계, 태그를 직접 조정합니다.",
            "출력 화면에서는 narrative 결과가 실제 UI 안에서 보입니다.",
            "이미지는 설명 대상이 아니라 생성 조건과 맥락의 일부입니다.",
        ),
    ),
    SlideSpec(
        slide_no=3,
        title="개인 맥락이 곧 포트폴리오의 신뢰도다",
        subtitle="프로필 요약과 인물 이미지를 분리하지 않고 한 장에서 연결합니다.",
        primary_image="image5.png",
        support_images=("image6.jpeg",),
        primary_role="profile_summary",
        visual_type="profile_summary_plus_portrait",
        supporting_text_goal="정리된 이력과 실제 인물 이미지가 함께 신뢰도를 만든다는 점을 강조한다.",
        layout_role="wide_profile_with_portrait_sidebar",
        crop_required=False,
        bullets=(
            "학력, 경력, 관심사를 설명문이 아니라 실제 자료로 제시합니다.",
            "개발자 정체성을 개인 맥락과 연결해 해석 기준을 만듭니다.",
        ),
    ),
    SlideSpec(
        slide_no=4,
        title="구현 깊이는 문제와 해법을 한 화면에 둔다",
        subtitle="문제 설명과 Python 구현을 동시에 보이게 해서 해결 흔적 자체를 자산으로 사용합니다.",
        primary_image="image22.png",
        support_images=(),
        primary_role="problem_and_code_proof",
        visual_type="code_and_problem_screen",
        supporting_text_goal="문제 해결형 개발자라는 점을 시각 증거로 보여준다.",
        layout_role="wide_problem_code_canvas",
        crop_required=False,
        bullets=(
            "문제 정의를 원문 화면으로 보여줍니다.",
            "해결 코드는 구현 수준의 디테일을 그대로 남깁니다.",
        ),
    ),
    SlideSpec(
        slide_no=5,
        title="실행 이력은 완료된 프로젝트 표로 확인한다",
        subtitle="프로젝트의 수, 역할, 진행 상태를 한 장의 evidence table로 정리합니다.",
        primary_image="image23.png",
        support_images=(),
        primary_role="portfolio_evidence_table",
        visual_type="table_evidence",
        supporting_text_goal="실행 범위와 지속성을 포트폴리오 증거 표로 보여준다.",
        layout_role="left_table_with_right_evidence_notes",
        crop_required=False,
        bullets=(
            "프로젝트를 역할, 기간, 상태 기준으로 비교합니다.",
            "포트폴리오가 실행 이력을 가진 문서임을 보여줍니다.",
            "다음 deck regeneration 작업과 바로 연결됩니다.",
        ),
    ),
    SlideSpec(
        slide_no=6,
        title="적용형 AI 워크플로는 화면 안에서 끝까지 닫힌다",
        subtitle="업로드, 조건 선택, 생성 결과가 한 서비스 화면 안에서 닫힌 구조를 보여줍니다.",
        primary_image="image29.png",
        support_images=(),
        primary_role="applied_ai_workflow_ui",
        visual_type="workflow_ui",
        supporting_text_goal="실험 결과를 실제 사용 흐름으로 이어 붙일 수 있다는 점을 강조한다.",
        layout_role="right_tall_ui_with_left_process_story",
        crop_required=False,
        bullets=(
            "이미지 업로드부터 생성 결과까지 한 인터페이스 안에서 이어집니다.",
            "멀티모달 입력은 noise가 아니라 form-preserving context로 다뤄집니다.",
            "다음 단계에서는 caption/eval/regeneration skill family로 확장할 수 있습니다.",
        ),
    ),
]


def load_manifest_index() -> dict[str, dict]:
    data = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    return {item["file"]: item for item in data["exported_images"]}


def rgb_to_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullet_block(slide, x, y, w, h, bullets, *, bg, text_color, size=15, step=0.72):
    add_rect(slide, x, y, w, h, bg)
    for idx, bullet in enumerate(bullets):
        add_textbox(
            slide,
            x + 0.22,
            y + 0.18 + idx * step,
            w - 0.44,
            0.58,
            f"• {bullet}",
            size=size,
            color=text_color,
            valign=MSO_ANCHOR.MIDDLE,
        )


def add_tag(slide, x, y, text, *, bg, fg):
    width = max(1.1, 0.14 * len(text) + 0.34)
    add_rect(slide, x, y, width, 0.46, bg)
    add_textbox(
        slide,
        x + 0.08,
        y + 0.03,
        width - 0.16,
        0.36,
        text,
        size=11,
        color=fg,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return width


def add_picture_contained(slide, image_path: Path, x, y, w, h, *, frame_fill=None, frame_line=None):
    if frame_fill is not None:
        add_rect(slide, x, y, w, h, frame_fill, line=frame_line)
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    image_ratio = img_w / img_h
    box_ratio = w / h
    if image_ratio > box_ratio:
        target_w = w
        target_h = w / image_ratio
    else:
        target_h = h
        target_w = h * image_ratio
    target_x = x + (w - target_w) / 2
    target_y = y + (h - target_h) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(target_x),
        Inches(target_y),
        width=Inches(target_w),
        height=Inches(target_h),
    )


def base_slide(prs: Presentation, *, dark: bool):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG if dark else LIGHT_BG
    return slide


def add_footer(slide, text, *, dark):
    add_textbox(
        slide,
        0.7,
        7.0,
        11.9,
        0.24,
        text,
        size=10,
        color=SOFT_TEXT if dark else RGBColor(0x76, 0x78, 0x7B),
        valign=MSO_ANCHOR.MIDDLE,
    )


def render_slide(prs: Presentation, spec: SlideSpec, image_index: dict[str, dict]):
    dark = spec.slide_no in {1, 6}
    slide = base_slide(prs, dark=dark)
    title_color = LIGHT_TEXT if dark else DARK_TEXT
    subtitle_color = RGBColor(0xD9, 0xE2, 0xE6) if dark else RGBColor(0x5A, 0x5F, 0x66)
    primary_image_path = Path(image_index[spec.primary_image]["output_path"])

    if spec.slide_no == 1:
        add_tag(slide, 0.75, 0.62, "LEAN 02_1 PORTFOLIO", bg=TEAL, fg=DARK_BG)
        add_tag(slide, 2.8, 0.62, "SYSTEM-FIRST", bg=SAND, fg=DARK_BG)
        add_textbox(slide, 0.75, 1.1, 5.1, 1.0, spec.title, size=28, color=title_color, bold=True)
        add_textbox(slide, 0.78, 2.1, 5.0, 1.3, spec.subtitle, size=16, color=subtitle_color)
        add_bullet_block(slide, 0.75, 3.55, 4.8, 2.55, spec.bullets, bg=RGBColor(0x2A, 0x35, 0x40), text_color=LIGHT_TEXT)
        add_picture_contained(slide, primary_image_path, 6.05, 0.9, 6.45, 5.9, frame_fill=RGBColor(0x23, 0x2D, 0x37), frame_line=RGBColor(0x3A, 0x4D, 0x59))
        add_footer(slide, "source: 02_1/image27.png • architecture hero", dark=True)
        return

    if spec.slide_no == 2:
        add_textbox(slide, 0.8, 0.5, 7.0, 0.9, spec.title, size=23, color=title_color, bold=True)
        add_textbox(slide, 0.82, 1.2, 6.8, 0.42, spec.subtitle, size=12.5, color=subtitle_color)
        add_picture_contained(slide, primary_image_path, 0.82, 1.65, 3.35, 4.95, frame_fill=RGBColor(0xEF, 0xE9, 0xDE), frame_line=RGBColor(0xD5, 0xC7, 0xB3))
        add_picture_contained(slide, Path(image_index[spec.support_images[0]]["output_path"]), 4.35, 1.65, 3.15, 4.95, frame_fill=RGBColor(0xEF, 0xE9, 0xDE), frame_line=RGBColor(0xD5, 0xC7, 0xB3))
        add_bullet_block(slide, 7.8, 1.85, 4.7, 3.05, spec.bullets, bg=MUTED_BG, text_color=DARK_TEXT, size=13.5, step=0.78)
        add_textbox(slide, 7.95, 5.18, 4.35, 0.95, "입력과 결과를 따로 나열하지 않고\n하나의 사용자 흐름으로 읽히게 배치합니다.", size=12.5, color=DARK_TEXT)
        add_footer(slide, "source: 02_1/image4.png + image37.png • paired product flow", dark=False)
        return

    if spec.slide_no == 3:
        add_textbox(slide, 0.8, 0.55, 6.6, 0.6, spec.title, size=24, color=title_color, bold=True)
        add_textbox(slide, 0.82, 1.08, 6.5, 0.55, spec.subtitle, size=13.5, color=subtitle_color)
        add_picture_contained(slide, primary_image_path, 0.82, 1.65, 7.1, 4.95, frame_fill=RGBColor(0xEF, 0xE9, 0xDE), frame_line=RGBColor(0xD5, 0xC7, 0xB3))
        add_picture_contained(slide, Path(image_index[spec.support_images[0]]["output_path"]), 8.45, 1.78, 3.05, 3.45, frame_fill=RGBColor(0xEF, 0xE9, 0xDE), frame_line=RGBColor(0xD5, 0xC7, 0xB3))
        add_bullet_block(slide, 7.95, 5.0, 4.65, 1.82, spec.bullets, bg=MUTED_BG, text_color=DARK_TEXT, size=12.5, step=0.5)
        add_footer(slide, "source: 02_1/image5.png + image6.jpeg • profile credibility pair", dark=False)
        return

    if spec.slide_no == 4:
        add_textbox(slide, 0.8, 0.55, 7.3, 0.6, spec.title, size=24, color=title_color, bold=True)
        add_textbox(slide, 0.82, 1.08, 7.2, 0.55, spec.subtitle, size=13.5, color=subtitle_color)
        add_picture_contained(slide, primary_image_path, 0.82, 1.55, 11.8, 4.35, frame_fill=RGBColor(0xEF, 0xE9, 0xDE), frame_line=RGBColor(0xD5, 0xC7, 0xB3))
        add_bullet_block(slide, 0.82, 5.98, 11.8, 0.92, spec.bullets, bg=MUTED_BG, text_color=DARK_TEXT, size=13.5, step=0.24)
        add_footer(slide, "source: 02_1/image22.png • problem and code proof", dark=False)
        return

    if spec.slide_no == 5:
        add_textbox(slide, 0.8, 0.55, 7.3, 0.6, spec.title, size=24, color=title_color, bold=True)
        add_textbox(slide, 0.82, 1.08, 7.0, 0.55, spec.subtitle, size=13.5, color=subtitle_color)
        add_picture_contained(slide, primary_image_path, 0.82, 1.65, 7.8, 4.85, frame_fill=RGBColor(0xEF, 0xE9, 0xDE), frame_line=RGBColor(0xD5, 0xC7, 0xB3))
        add_bullet_block(slide, 8.95, 1.85, 3.55, 2.75, spec.bullets, bg=MUTED_BG, text_color=DARK_TEXT, size=12.8, step=0.62)
        add_textbox(slide, 9.05, 5.3, 3.25, 0.95, "표 이미지는 form 자체가 의미를 가지므로\n요약 대신 실제 evidence block으로 남깁니다.", size=12.5, color=DARK_TEXT)
        add_footer(slide, "source: 02_1/image23.png • portfolio evidence table", dark=False)
        return

    add_tag(slide, 0.75, 0.62, "APPLIED AI WORKFLOW", bg=TEAL, fg=DARK_BG)
    add_textbox(slide, 0.75, 1.12, 5.45, 0.85, spec.title, size=27, color=title_color, bold=True)
    add_textbox(slide, 0.78, 2.05, 5.0, 1.0, spec.subtitle, size=15.5, color=subtitle_color)
    add_bullet_block(slide, 0.75, 3.2, 5.0, 2.55, spec.bullets, bg=RGBColor(0x2A, 0x35, 0x40), text_color=LIGHT_TEXT)
    add_picture_contained(slide, primary_image_path, 6.25, 0.95, 6.0, 5.95, frame_fill=RGBColor(0x23, 0x2D, 0x37), frame_line=RGBColor(0x3A, 0x4D, 0x59))
    add_footer(slide, "source: 02_1/image29.png • applied AI workflow UI", dark=True)


def build_role_matrix(image_index: dict[str, dict]) -> list[dict]:
    rows = []
    for spec in SLIDE_SPECS:
        rows.append(
            {
                "slide_no": spec.slide_no,
                "slide_title": spec.title,
                "visual_type": spec.visual_type,
                "image_filename": spec.primary_image,
                "support_images": list(spec.support_images),
                "portfolio_role": spec.primary_role,
                "source_slide_numbers": image_index[spec.primary_image]["slide_usages"],
                "supporting_text_goal": spec.supporting_text_goal,
                "layout_role": spec.layout_role,
                "crop_required": spec.crop_required,
            }
        )
    return rows


def build_presentation(slides: list[SlideSpec], image_index: dict[str, dict], output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_SIZE[0])
    prs.slide_height = Inches(SLIDE_SIZE[1])
    for spec in slides:
        render_slide(prs, spec, image_index)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main():
    image_index = load_manifest_index()
    build_presentation(SLIDE_SPECS, image_index, DECK_PATH)
    for spec in SLIDE_SPECS:
        build_presentation([spec], image_index, RENDER_SOURCE_DIR / f"slide-{spec.slide_no:02d}-source.pptx")
    ROLE_MATRIX_PATH.write_text(
        json.dumps(build_role_matrix(image_index), ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    print(
        json.dumps(
            {
                "deck": str(DECK_PATH),
                "render_sources": [str(RENDER_SOURCE_DIR / f"slide-{spec.slide_no:02d}-source.pptx") for spec in SLIDE_SPECS],
                "role_matrix": str(ROLE_MATRIX_PATH),
                "authoring_skill": "/Users/jaehyuntak/.codex/skills/pptx/SKILL.md",
                "copy_clarity_skill": "/Users/jaehyuntak/.claude/skills/semantic-clarity-enhanced/SKILL.md",
                "palette": {
                    "dark_bg": rgb_to_hex(DARK_BG),
                    "light_bg": rgb_to_hex(LIGHT_BG),
                    "accent_teal": rgb_to_hex(TEAL),
                    "accent_sand": rgb_to_hex(SAND),
                },
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
